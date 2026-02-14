from fastapi import UploadFile
import pdfplumber
import docx
import logging
import numpy as np
import cv2
from pdf2image import convert_from_bytes
from pptx import Presentation 
import io

logger = logging.getLogger(__name__)

MIN_TEXT_LENGTH = 50 

def extract_text(file: UploadFile) -> str:

    file.file.seek(0)
    filename = file.filename.lower()

    try:
        if filename.endswith(".pdf"):
            return _extract_pdf_text(file)

        elif filename.endswith(".docx"):
            return _extract_docx(file)

        elif filename.endswith(".pptx") or filename.endswith(".ppt"):
            # <--- NEW HANDLER FOR PPT
            return _extract_pptx(file)

        elif filename.endswith(".txt"):
            return _extract_txt(file)

        else:
            # Fallback for unsupported images or other types
            return "OCR is currently disabled. Please upload a text-based PDF, DOCX, PPTX, or TXT file."

    except Exception as e:
        logger.exception("Text extraction failed")
        # Return empty string or error message so the app doesn't crash
        return "" 


def _extract_pdf_text(file: UploadFile) -> str:
    text_blocks = []
    with pdfplumber.open(file.file) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_blocks.append(page_text)
    return "\n".join(text_blocks)

def _extract_docx(file: UploadFile) -> str:

    file_content = file.file.read()
    doc = docx.Document(io.BytesIO(file_content))
    return "\n".join(p.text for p in doc.paragraphs)

def _extract_pptx(file: UploadFile) -> str:

    file_content = file.file.read()
    prs = Presentation(io.BytesIO(file_content))
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # 1. Join runs with a space to prevent "greedyield"
                    para_text = " ".join([run.text for run in paragraph.runs])
                    if para_text.strip():
                        text_runs.append(para_text)

    return "\n".join(text_runs)

def _extract_txt(file: UploadFile) -> str:
    return file.file.read().decode("utf-8", errors="ignore")


def _extract_pdf_ocr(file: UploadFile) -> str:
    # Placeholder for your OCR logic
    return ""

def _extract_image_ocr(file: UploadFile) -> str:
    # Placeholder for your OCR logic
    return ""